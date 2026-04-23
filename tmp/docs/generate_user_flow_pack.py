#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from workflow_pack_data import (
    ASSET_ROOT,
    DECK_OUTPUT_DIR,
    GENERATED_ON,
    INVENTORY_DEPRECATED,
    INVENTORY_MISSING,
    ROLE_MAP,
    SOURCE_MANUAL_DIR,
    WORKFLOW_GROUPS,
    WORKFLOWS,
    select_workflows,
)


REPO_ROOT = Path(__file__).resolve().parents[2]
MANUAL_DIR = REPO_ROOT / SOURCE_MANUAL_DIR
ASSET_DIR = REPO_ROOT / ASSET_ROOT
DECK_DIR = REPO_ROOT / DECK_OUTPUT_DIR

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(11, 31, 51)
TEAL = RGBColor(20, 184, 166)
ORANGE = RGBColor(249, 115, 22)
SURFACE = RGBColor(247, 248, 250)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(16, 42, 67)
MUTED = RGBColor(82, 102, 122)
SOFT = RGBColor(226, 232, 240)


def ensure_dirs() -> None:
    MANUAL_DIR.mkdir(parents=True, exist_ok=True)
    DECK_DIR.mkdir(parents=True, exist_ok=True)


def manual_path(workflow: dict) -> Path:
    return MANUAL_DIR / f"{workflow['order']:02d}-{workflow['slug']}.md"


def deck_path(workflow: dict) -> Path:
    return DECK_DIR / f"{workflow['order']:02d}-{workflow['slug']}.pptx"


def asset_rel_path(workflow: dict, filename: str) -> str:
    return f"{ASSET_ROOT}/{workflow['slug']}/{filename}"


def asset_abs_path(workflow: dict, filename: str) -> Path:
    return ASSET_DIR / workflow["slug"] / filename


def write_manual_index() -> None:
    lines = [
        "# Product User-Flow Manuals",
        "",
        f"Generated from the live demo-aligned workflow source on {GENERATED_ON}.",
        "",
        "## Workflow Inventory",
        "",
        "| No. | Workflow | Review status | Delta note | Primary user | Source manual |",
        "| --- | --- | --- | --- | --- | --- |",
    ]
    for wf in WORKFLOWS:
        lines.append(
            f"| {wf['order']:02d} | {wf['title']} | {wf.get('inventory_status', '✅ up-to-date')} | "
            f"{wf.get('inventory_note', 'No material drift found during the current review.')} | "
            f"{wf['primary_user']} | `{manual_path(wf).relative_to(REPO_ROOT)}` |"
        )
    lines.extend(
        [
            "",
            "## Deprecated or Compatibility-Only Items",
            "",
        ]
    )
    if INVENTORY_DEPRECATED:
        for item in INVENTORY_DEPRECATED:
            lines.append(f"- {item}")
    else:
        lines.append("- None identified in the current review.")

    lines.extend(
        [
            "",
            "## Missing Workflows Identified",
            "",
        ]
    )
    if INVENTORY_MISSING:
        for item in INVENTORY_MISSING:
            lines.append(f"- {item}")
    else:
        lines.append("- No additional top-level workflows were discovered beyond the existing pack.")

    lines.extend(
        [
            "",
            "## Asset Root",
            "",
            f"`{ASSET_ROOT}/`",
            "",
            "## Deck Output",
            "",
            f"`{DECK_OUTPUT_DIR}/`",
            "",
        ]
    )
    (MANUAL_DIR / "README.md").write_text("\n".join(lines), encoding="utf-8")


def build_manual(workflow: dict) -> str:
    lines = [
        f"# {workflow['title']}",
        "",
        "## 1. Title",
        "",
        workflow["title"],
        "",
        "## 2. Document Purpose",
        "",
        workflow["document_purpose"],
        "",
        "## 3. Primary User",
        "",
        workflow["primary_user"],
        "",
        "## 4. Entry Point",
        "",
        workflow["entry_point"],
        "",
        "## 5. Workflow Summary",
        "",
    ]

    for item in workflow["workflow_summary"]:
        lines.append(f"- {item}")

    if workflow.get("extra_markdown"):
        lines.extend(["", workflow["extra_markdown"].rstrip(), ""])

    lines.extend(["## 6. Step-By-Step Instructions", ""])

    for step in workflow["steps"]:
        lines.extend(
            [
                f"### Step {step['number']}. {step['title']}",
                "",
                f"- What the user does: {step['user_does']}",
                f"- What the user sees: {step['user_sees']}",
                f"- Why the step matters: {step['why_it_matters']}",
                f"- Expected result: {step['expected_result']}",
                f"- Common issues / trainer notes: {step['trainer_notes']}",
                "- Screenshot placeholder:",
                f"  Suggested file path: `{asset_rel_path(workflow, step['screenshot_file'])}`",
                f"  Screenshot caption: {step['screenshot_caption']}",
                f"  What the screenshot should show: {step['screenshot_focus']}",
                "",
            ]
        )

    lines.extend(["## 7. Success Criteria", ""])
    for item in workflow["success_criteria"]:
        lines.append(f"- {item}")

    lines.extend(["", "## 8. Related Documents", ""])
    for item in workflow["related_documents"]:
        lines.append(f"- `{item}`")

    lines.extend(
        [
            "",
            "## 9. Status",
            "",
            workflow["status"],
            "",
        ]
    )
    return "\n".join(lines)


def write_manuals() -> None:
    for wf in WORKFLOWS:
        manual_path(wf).write_text(build_manual(wf), encoding="utf-8")


def write_selected_manuals(selected_workflows: list[dict]) -> None:
    for wf in selected_workflows:
        manual_path(wf).write_text(build_manual(wf), encoding="utf-8")


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line_color=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line_color or color
    return shape


def add_textbox(slide, left, top, width, height, text, *, font_size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, font_size=20, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return box


def add_step_sections(slide, left, top, width, height, step: dict) -> None:
    sections = [
        ("User action", step["user_does"], TEAL, TEXT),
        ("What appears on screen", step["user_sees"], TEAL, TEXT),
        ("Why it matters", step["why_it_matters"], TEAL, TEXT),
        ("Expected result", step["expected_result"], TEAL, TEXT),
        ("Trainer note", step["trainer_notes"], ORANGE, MUTED),
    ]

    total_chars = sum(len(body) for _, body, _, _ in sections)
    body_font = 16
    if total_chars > 720:
        body_font = 11
    elif total_chars > 580:
        body_font = 12
    elif total_chars > 430:
        body_font = 13
    elif total_chars > 320:
        body_font = 14
    label_font = 12
    trainer_font = max(body_font - 1, 11)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    first = True
    for label, body, label_color, body_color in sections:
        p_label = tf.paragraphs[0] if first else tf.add_paragraph()
        p_label.text = label
        p_label.alignment = PP_ALIGN.LEFT
        p_label.space_before = Pt(0)
        p_label.space_after = Pt(2)
        for run in p_label.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(label_font)
            run.font.bold = True
            run.font.color.rgb = label_color

        p_body = tf.add_paragraph()
        p_body.text = body
        p_body.alignment = PP_ALIGN.LEFT
        p_body.space_before = Pt(0)
        p_body.space_after = Pt(8 if label != "Trainer note" else 0)
        p_body.line_spacing = 1.1
        for run in p_body.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(trainer_font if label == "Trainer note" else body_font)
            run.font.color.rgb = body_color

        first = False


def add_title_band(slide, title: str, subtitle: str | None = None) -> None:
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.65), NAVY)
    add_rect(slide, Inches(0), Inches(0.65), Inches(1.35), Inches(0.12), ORANGE)
    if len(title) > 42:
        title_font = 18
    elif len(title) > 32:
        title_font = 21
    else:
        title_font = 26
    add_textbox(slide, Inches(0.55), Inches(0.12), Inches(10.2), Inches(0.32), title, font_size=title_font, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(9.3), Inches(0.2), Inches(3.4), Inches(0.25), subtitle, font_size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def add_footer(slide, text: str) -> None:
    add_textbox(slide, Inches(0.55), Inches(7.0), Inches(12.0), Inches(0.2), text, font_size=10, color=MUTED)


def fit_image(slide, path: Path, left, top, width, height) -> None:
    if not path.exists():
        add_rect(slide, left, top, width, height, SOFT, radius=True)
        add_textbox(slide, left + Inches(0.3), top + Inches(0.35), width - Inches(0.6), height - Inches(0.7), f"Screenshot pending\n{path.name}", font_size=18, color=MUTED, align=PP_ALIGN.CENTER)
        return

    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_cover_slide(prs: Presentation, workflow: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_rect(slide, Inches(0.55), Inches(0.8), Inches(0.18), Inches(5.6), ORANGE)
    add_textbox(slide, Inches(0.95), Inches(0.95), Inches(7.35), Inches(1.55), workflow["title"], font_size=26, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.95), Inches(2.35), Inches(7.15), Inches(1.45), workflow["document_purpose"], font_size=16, color=WHITE)
    add_rect(slide, Inches(0.95), Inches(4.6), Inches(2.6), Inches(0.5), TEAL, radius=True)
    add_textbox(slide, Inches(1.2), Inches(4.72), Inches(2.2), Inches(0.2), workflow["deck_group"], font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.95), Inches(5.45), Inches(6.95), Inches(0.72), f"Primary user: {workflow['primary_user']}", font_size=15, color=WHITE)
    add_textbox(slide, Inches(0.95), Inches(6.28), Inches(3.5), Inches(0.25), f"Generated: {GENERATED_ON}", font_size=11, color=WHITE)

    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.05), Inches(3.85), Inches(4.9))
    hero.fill.solid()
    hero.fill.fore_color.rgb = WHITE
    hero.line.color.rgb = WHITE
    tf = hero.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "Training Pack"
    p1.alignment = PP_ALIGN.CENTER
    for run in p1.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = textwrap.fill(workflow["entry_point"], width=28)
    p2.alignment = PP_ALIGN.CENTER
    for run in p2.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(13)
        run.font.color.rgb = MUTED


def add_overview_slide(prs: Presentation, workflow: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SURFACE)
    add_title_band(slide, "Workflow Overview", workflow["title"])

    cards = [
        ("Primary User", workflow["primary_user"]),
        ("Entry Point", workflow["entry_point"]),
        ("Purpose", workflow["document_purpose"]),
    ]
    top = Inches(1.05)
    lefts = [Inches(0.55), Inches(4.55), Inches(8.55)]
    for (label, value), left in zip(cards, lefts):
        add_rect(slide, left, top, Inches(3.6), Inches(1.45), WHITE, line_color=SOFT, radius=True)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.15), Inches(3.2), Inches(0.25), label, font_size=12, color=TEAL, bold=True)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.45), Inches(3.15), Inches(0.8), value, font_size=14, color=TEXT)

    add_rect(slide, Inches(0.55), Inches(2.85), Inches(12.2), Inches(3.55), WHITE, line_color=SOFT, radius=True)
    add_textbox(slide, Inches(0.8), Inches(3.05), Inches(2.2), Inches(0.25), "What to cover", font_size=13, color=TEAL, bold=True)
    add_bullets(slide, Inches(0.8), Inches(3.38), Inches(11.6), Inches(2.7), workflow["workflow_summary"], font_size=18)
    add_footer(slide, f"Source manual: {manual_path(workflow).relative_to(REPO_ROOT)}")


def add_role_map_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SURFACE)
    add_title_band(slide, "Role Map", "Platform Overview")

    top = Inches(1.05)
    left = Inches(0.55)
    width = Inches(12.15)
    height = Inches(0.82)
    for idx, item in enumerate(ROLE_MAP):
        y = top + Inches(idx * 0.92)
        add_rect(slide, left, y, width, height, WHITE, line_color=SOFT, radius=True)
        add_textbox(slide, left + Inches(0.2), y + Inches(0.14), Inches(2.1), Inches(0.22), item["role"], font_size=15, color=TEAL, bold=True)
        add_textbox(slide, left + Inches(2.45), y + Inches(0.12), Inches(3.7), Inches(0.5), item["goal"], font_size=12, color=TEXT)
        add_textbox(slide, left + Inches(6.35), y + Inches(0.12), Inches(2.65), Inches(0.5), item["starts_at"], font_size=11, color=MUTED)
        add_textbox(slide, left + Inches(9.3), y + Inches(0.12), Inches(2.5), Inches(0.5), item["hands_off_to"], font_size=11, color=MUTED)
    add_footer(slide, "Cross-role handoff map for the seeded training campaign.")


def add_step_slide(prs: Presentation, workflow: dict, step: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SURFACE)
    add_title_band(slide, f"Step {step['number']}: {step['title']}", workflow["title"])

    add_rect(slide, Inches(0.55), Inches(1.0), Inches(5.0), Inches(5.7), WHITE, line_color=SOFT, radius=True)
    add_step_sections(slide, Inches(0.82), Inches(1.2), Inches(4.3), Inches(5.2), step)

    add_rect(slide, Inches(5.8), Inches(1.0), Inches(6.95), Inches(5.7), WHITE, line_color=SOFT, radius=True)
    add_textbox(slide, Inches(6.1), Inches(1.2), Inches(6.2), Inches(0.25), step["screenshot_caption"], font_size=13, color=TEAL, bold=True)
    fit_image(slide, asset_abs_path(workflow, step["screenshot_file"]), Inches(6.08), Inches(1.55), Inches(6.35), Inches(4.55))
    add_textbox(slide, Inches(6.1), Inches(6.2), Inches(6.15), Inches(0.35), step["screenshot_focus"], font_size=11, color=MUTED)
    add_footer(slide, f"Screenshot asset: {asset_rel_path(workflow, step['screenshot_file'])}")


def add_tips_slide(prs: Presentation, workflow: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SURFACE)
    add_title_band(slide, "Trainer Tips", workflow["title"])
    add_rect(slide, Inches(0.7), Inches(1.3), Inches(12.0), Inches(4.8), WHITE, line_color=SOFT, radius=True)
    add_textbox(slide, Inches(1.0), Inches(1.6), Inches(2.5), Inches(0.3), "Facilitation notes", font_size=14, color=TEAL, bold=True)
    add_bullets(slide, Inches(1.0), Inches(2.0), Inches(11.2), Inches(3.3), workflow["tips"], font_size=22)
    add_textbox(slide, Inches(1.0), Inches(5.75), Inches(11.0), Inches(0.4), "Related documents: " + ", ".join(workflow["related_documents"]), font_size=11, color=MUTED)
    add_footer(slide, workflow["status"])


def add_close_slide(prs: Presentation, workflow: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(5.8), Inches(0.55), "What success looks like", font_size=28, color=WHITE, bold=True)
    add_bullets(slide, Inches(0.85), Inches(1.7), Inches(5.7), Inches(3.5), workflow["success_criteria"], font_size=20, color=WHITE)
    add_rect(slide, Inches(7.1), Inches(1.25), Inches(5.4), Inches(3.95), TEAL, radius=True)
    add_textbox(slide, Inches(7.45), Inches(1.6), Inches(4.6), Inches(0.3), "Use next", font_size=16, color=WHITE, bold=True)
    add_bullets(slide, Inches(7.45), Inches(2.0), Inches(4.45), Inches(2.4), [f"Review {doc}" for doc in workflow["related_documents"][:3]], font_size=18, color=WHITE)
    add_textbox(slide, Inches(0.85), Inches(6.2), Inches(11.7), Inches(0.35), workflow["status"], font_size=12, color=WHITE)


def generate_workflow_deck(workflow: dict) -> None:
    prs = new_prs()
    add_cover_slide(prs, workflow)
    add_overview_slide(prs, workflow)
    if workflow["order"] == 1:
        add_role_map_slide(prs)
    for step in workflow["steps"]:
        add_step_slide(prs, workflow, step)
    add_tips_slide(prs, workflow)
    add_close_slide(prs, workflow)
    prs.save(deck_path(workflow))


def generate_template_deck() -> None:
    prs = new_prs()

    sample = {
        "order": 99,
        "slug": "user-flow-template",
        "title": "Reusable Training Deck Template",
        "deck_group": "Template",
        "document_purpose": "Template deck showing the standard slide types used across the product training pack.",
        "primary_user": "Documentation and enablement teams.",
        "entry_point": "Use as the styling baseline for future workflow decks.",
        "workflow_summary": [
            "Cover slide for the workflow title and audience.",
            "Overview slide for purpose, user, and entry point.",
            "Step detail slide with screenshot and trainer notes.",
            "Trainer tips and close-out slides for delivery.",
        ],
        "related_documents": ["docs/product-user-flows/README.md"],
        "status": f"Generated on {GENERATED_ON}.",
        "tips": [
            "Keep screenshots large and readable.",
            "Prefer one screen per step rather than dense multi-screen collages.",
        ],
        "success_criteria": [
            "The deck feels presentation-ready and editable.",
            "Layout consistency matches the workflow decks in this pack.",
        ],
    }

    add_cover_slide(prs, sample)
    add_overview_slide(prs, sample)

    step_slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(step_slide, SURFACE)
    add_title_band(step_slide, "Wide Screenshot Layout", "Template")
    add_rect(step_slide, Inches(0.6), Inches(1.05), Inches(12.1), Inches(5.7), WHITE, line_color=SOFT, radius=True)
    add_textbox(step_slide, Inches(0.95), Inches(1.35), Inches(11.0), Inches(0.28), "Use this layout when a single page screenshot carries most of the teaching value.", font_size=16, color=TEXT)
    add_rect(step_slide, Inches(0.95), Inches(1.95), Inches(11.35), Inches(4.2), SOFT, radius=True)
    add_textbox(step_slide, Inches(1.2), Inches(3.4), Inches(10.6), Inches(0.5), "Insert full-width screenshot here", font_size=24, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

    decision = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(decision, SURFACE)
    add_title_band(decision, "Comparison / Decision Slide", "Template")
    add_rect(decision, Inches(0.8), Inches(1.45), Inches(5.7), Inches(4.6), WHITE, line_color=SOFT, radius=True)
    add_rect(decision, Inches(6.8), Inches(1.45), Inches(5.7), Inches(4.6), WHITE, line_color=SOFT, radius=True)
    add_textbox(decision, Inches(1.1), Inches(1.75), Inches(5.0), Inches(0.3), "Option A", font_size=20, color=TEAL, bold=True)
    add_bullets(decision, Inches(1.1), Inches(2.2), Inches(4.9), Inches(3.1), ["When the workflow is role-specific", "When one screen is the primary teaching artifact"], font_size=18)
    add_textbox(decision, Inches(7.1), Inches(1.75), Inches(5.0), Inches(0.3), "Option B", font_size=20, color=ORANGE, bold=True)
    add_bullets(decision, Inches(7.1), Inches(2.2), Inches(4.9), Inches(3.1), ["When a deck needs a branch or exception path", "When trainers need a decision checklist"], font_size=18)

    add_tips_slide(prs, sample)
    add_close_slide(prs, sample)
    prs.save(DECK_DIR / "zz-user-flow-template.pptx")


def add_hyperlink_run(paragraph, text: str, address: str, color=TEAL, font_size: int = 18) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = color
    run.hyperlink.address = address


def add_index_entry(slide, workflow: dict, left, top, width) -> None:
    box = slide.shapes.add_textbox(left, top, width, Inches(0.62))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    title = f"{workflow['order']:02d}. {workflow['title']}"
    title_font = 12 if len(title) > 46 else 14
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    add_hyperlink_run(p1, title, deck_path(workflow).name, color=RGBColor(29, 78, 216), font_size=title_font)

    p2 = tf.add_paragraph()
    p2.text = workflow["primary_user"]
    p2.alignment = PP_ALIGN.LEFT
    for run in p2.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(11)
        run.font.color.rgb = MUTED


def generate_master_index() -> None:
    prs = new_prs()

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(cover, NAVY)
    add_textbox(cover, Inches(0.9), Inches(1.05), Inches(7.6), Inches(0.8), "User Flow Training Index", font_size=28, color=WHITE, bold=True)
    add_textbox(cover, Inches(0.9), Inches(2.0), Inches(6.8), Inches(0.9), "Open this deck first. Each linked item opens the matching workflow deck in the same shared folder.", font_size=17, color=WHITE)
    add_rect(cover, Inches(0.9), Inches(3.25), Inches(3.3), Inches(0.72), ORANGE, radius=True)
    add_textbox(cover, Inches(1.08), Inches(3.46), Inches(2.95), Inches(0.28), f"{len(WORKFLOWS)} workflow decks + template", font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    index_slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(index_slide, SURFACE)
    add_title_band(index_slide, "Linked Deck Catalogue", "Open the sibling PPTX files from this slide")

    columns = [
        (WORKFLOW_GROUPS[:3], Inches(0.75), Inches(1.0), Inches(5.65)),
        (WORKFLOW_GROUPS[3:], Inches(6.75), Inches(1.0), Inches(5.75)),
    ]
    for groups, left, top, width in columns:
        y = top
        for group_name, order_list in groups:
            add_textbox(index_slide, left, y, width, Inches(0.24), group_name, font_size=15, color=TEAL, bold=True)
            y += Inches(0.32)
            for order in order_list:
                workflow = next(wf for wf in WORKFLOWS if wf["order"] == order)
                add_index_entry(index_slide, workflow, left + Inches(0.2), y, width - Inches(0.2))
                y += Inches(0.62)
            y += Inches(0.08)

    add_footer(index_slide, f"All links use sibling file names so the folder can be shared as one package. Generated {GENERATED_ON}.")

    prs.save(DECK_DIR / "00-user-flow-training-index.pptx")


def write_manifest() -> None:
    inventory = [
        {
            "order": wf["order"],
            "title": wf["title"],
            "source_manual": str(manual_path(wf).relative_to(REPO_ROOT)),
            "deck": str(deck_path(wf).relative_to(REPO_ROOT)),
        }
        for wf in WORKFLOWS
    ]
    (DECK_DIR / "user-flow-pack-manifest.json").write_text(json.dumps(inventory, indent=2), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate editable manuals and PPTX workflow decks.")
    parser.add_argument(
        "selectors",
        nargs="*",
        help="Optional workflow orders or slugs to refresh. Defaults to all workflows.",
    )
    parser.add_argument(
        "--workflows",
        dest="workflow_csv",
        help="Comma-separated workflow orders or slugs to refresh.",
    )
    parser.add_argument(
        "--include-template",
        action="store_true",
        help="Also regenerate the reusable template deck during a selective refresh.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    selectors = list(args.selectors)
    if args.workflow_csv:
        selectors.append(args.workflow_csv)
    selected_workflows = select_workflows(selectors)
    selective_refresh = bool(selectors)

    ensure_dirs()
    write_manual_index()
    if selective_refresh:
        write_selected_manuals(selected_workflows)
    else:
        write_manuals()
    for workflow in selected_workflows:
        generate_workflow_deck(workflow)
    if not selective_refresh or args.include_template:
        generate_template_deck()
    generate_master_index()
    write_manifest()
    selected_orders = ", ".join(f"{workflow['order']:02d}" for workflow in selected_workflows)
    print(f"Generated manuals and PPTX decks for workflows: {selected_orders}")


if __name__ == "__main__":
    main()
